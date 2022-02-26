import win32com.client
from os import walk
from pptx import Presentation # 라이브러리 
from pptx.util import Inches # 사진, 표등을 그리기 위해

prs = Presentation() # 파워포인트 객체 선언
file_dir = input('파일경로를 입력해주세요.:')
path,_,files = next(walk(file_dir))
outputFileName = path + '\\' + 'allSildesMerged.pptx'
img_path = []
num = 0
print(files)

for file2 in files:
    img_path.append(path + '\\' + file2)


for file in img_path:
    blank_slide_layout = prs.slide_layouts[6] # 6 : 제목/내용이 없는 '빈' 슬라이드
    slide = prs.slides.add_slide(blank_slide_layout)
    txBox = slide.shapes.add_textbox(1, 1, 5, 1) # 텍스트 박스 생성하기 
    tf = txBox.text_frame  #프레임 만들고
    tf.text = files[num]  #파일이름을 불러오기
    num += 1  # 파일이름이 리스트인 관계로 숫자를 올림

    left = top = Inches(1.2)
    width = Inches(7)
    height = Inches(6)
    # width, hegith가 없을 경우 원본 사이즈로 
    pic = slide.shapes.add_picture(file, left, top, width=width,height=height)
prs.save('add all slides.pptx')

