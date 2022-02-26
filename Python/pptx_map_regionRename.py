from pptx import Presentation
from pptx.util import Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
import pandas as pd
import tkinter
from tkinter import filedialog
from pptx.enum.text import PP_ALIGN

root = tkinter.Tk()
root.withdraw()
dir_path = filedialog.askopenfile(filetypes=[('엑셀','*.xlsx')])
print("\ndir_path : ", dir_path)

df = pd.read_excel(dir_path.name)

psr = Presentation('d:/공유용/4부문_지도(211201).pptx')
slide = psr.slides[0]
shape_lv1 = slide.shapes
shape_idx = []

shape_search = {}

for i in shape_lv1:
    if i.name == '4dept_all':
        for j in i.shapes:
            shape_idx.append(j)

for i, value  in enumerate(shape_idx):
    shape_search[value.name] = i
    a = df[df['지역']==value.name]['점포수']
    if a.empty == False:
        shape_idx[i].text = value.name + '\n(' + str(a.values[0]) +')'
        shape_idx[i].text_frame.paragraphs[0].font.name = '나눔고딕'
        shape_idx[i].text_frame.paragraphs[0].font.size = Pt(8)
        shape_idx[i].text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        shape_idx[i].text_frame.paragraphs[1].font.name = '나눔고딕'
        shape_idx[i].text_frame.paragraphs[1].font.size = Pt(6)  
        shape_idx[i].text_frame.paragraphs[1].alignment = PP_ALIGN.CENTER
        print(value.name + ': Success(' + shape_idx[i].text + ')')

# ddf = pd.DataFrame(list(shape_search.items()),columns=['1','2'])
# ddf.to_excel('text.xlsx')
psr.save('sss.pptx')

# if group_test == MSO_SHAPE_TYPE.TEXT_BOX:
#         new_name = i.text.split('(')[0]
#         if new_name != '' and 6 > int(len(new_name)):
#             new_name = new_name.strip('\n,\x0b')
#             i.name = new_name
#             # print(i.name)
#             # i.text = i.name + '\n(' + '5' + ')'
#             # i.text_frame.paragraphs[0].font.name = '나눔고딕'
#             # i.text_frame.paragraphs[0].font.size = Pt(9)
#             # i.text_frame.paragraphs[1].font.name = '나눔고딕'
#             # i.text_frame.paragraphs[1].font.size = Pt(9)
#             shape_idx.append(i)
#     if group_test == MSO_SHAPE_TYPE.GROUP:
#         for j in i.shapes:
#             group_test2 = j.shape_type
#             if group_test2 == MSO_SHAPE_TYPE.TEXT_BOX:
#                 new_name = j.text.split('(')[0]
#                 if new_name != '':
#                     j.name = new_name.strip('\n')
#                     # print(j.name)
#                     # j.text = j.name + '\n(' + '5' + ')'
#                     # j.text_frame.paragraphs[0].font.name = '나눔고딕'
#                     # j.text_frame.paragraphs[0].font.size = Pt(9)
#                     # j.text_frame.paragraphs[1].font.name = '나눔고딕'
#                     # j.text_frame.paragraphs[1].font.size = Pt(9)
#                     shape_idx.append(j)

#             if group_test2 == MSO_SHAPE_TYPE.GROUP:
#                 for c in j.shapes:
#                     group_test3 = c.shape_type
#                     if group_test3 == MSO_SHAPE_TYPE.TEXT_BOX:
#                         new_name = c.text.split('(')[0]
#                         if new_name != '':
#                             c.name = new_name.strip('\n')
#                             # print(c.name)
#                             # c.text = c.name + '\n(' + '5' + ')'
#                             # c.text_frame.paragraphs[0].font.name = '나눔고딕'
#                             # c.text_frame.paragraphs[0].font.size = Pt(9)
#                             # c.text_frame.paragraphs[1].font.name = '나눔고딕'
#                             # c.text_frame.paragraphs[1].font.size = Pt(9)
#                             shape_idx.append(c)

#                     if group_test3 == MSO_SHAPE_TYPE.GROUP:
#                         for d in c.shapes:
#                             group_test4 = d.shape_type
#                             if group_test4 == MSO_SHAPE_TYPE.TEXT_BOX:
#                                 new_name = d.text.split('(')[0]
#                                 if new_name != '':
#                                     d.name = new_name.strip('\n')
#                                     # print(d.name)
#                                     # d.text = d.name + '\n(' + '5' + ')'
#                                     # d.text_frame.paragraphs[0].font.name = '나눔고딕'
#                                     # d.text_frame.paragraphs[0].font.size = Pt(9)
#                                     # d.text_frame.paragraphs[1].font.name = '나눔고딕'
#                                     # d.text_frame.paragraphs[1].font.size = Pt(9)                                    
#                                     shape_idx.append(d)

#                             if group_test4 == MSO_SHAPE_TYPE.GROUP:
#                                 for b in d.shapes:
#                                     group_test5 = b.shape_type
#                                     if group_test5 == MSO_SHAPE_TYPE.TEXT_BOX:
#                                         new_name = b.text.split('(')[0]
#                                         if new_name != '':
#                                             b.name = new_name.strip('\n')
#                                             # print(b.name)
#                                             # b.text = b.name + '\n(' + '5' + ')'
#                                             # b.text_frame.paragraphs[0].font.name = '나눔고딕'
#                                             # b.text_frame.paragraphs[0].font.size = Pt(9)
#                                             # b.text_frame.paragraphs[1].font.name = '나눔고딕'
#                                             # b.text_frame.paragraphs[1].font.size = Pt(9)                                                     
#                                             shape_idx.append(d)
# shape_search = {}

# print(shape_idx)
# shape_index에 있는 오브젝트를 찾기 위해 shape_search로 인뎃스 값을 만듬

# for i, value  in enumerate(shape_idx):
#     shape_search[value.name] = i
#     a = df[df['지역']==value.name]['점포수']
#     shape_idx[i].text = value.name + '\n(' + str(a.values[0]) +')'
#     shape_idx[i].text_frame.paragraphs[0].font.name = '나눔고딕'
#     shape_idx[i].text_frame.paragraphs[0].font.size = Pt(9)
#     shape_idx[i].text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
#     shape_idx[i].text_frame.paragraphs[1].font.name = '나눔고딕'
#     shape_idx[i].text_frame.paragraphs[1].font.size = Pt(9)  
#     shape_idx[i].text_frame.paragraphs[1].alignment = PP_ALIGN.CENTER
#     print(i,j,shape_idx[i])

# psr.save('sss.pptx')